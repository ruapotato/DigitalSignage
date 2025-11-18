#!/usr/bin/env python3
"""
Digital Signage Management System
A simple Flask application for managing digital signage content across multiple TVs
"""

import os
import json
import uuid
from functools import wraps
from flask import Flask, render_template, request, redirect, url_for, session, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from pptx import Presentation
from PIL import Image
import io

app = Flask(__name__)
app.secret_key = os.urandom(24)  # Generate random secret key for sessions

# Configuration
SLIDES_DIR = 'slides'
ALLOWED_EXTENSIONS = {'pptx', 'jpg', 'jpeg', 'png', 'gif'}
MAX_TVS = 4
DEFAULT_SLIDE_DURATION = 5
TARGET_WIDTH = 1920
TARGET_HEIGHT = 1080

# Ensure slides directory exists
os.makedirs(SLIDES_DIR, exist_ok=True)


def load_credentials():
    """Load username and password from creds.txt"""
    try:
        with open('creds.txt', 'r') as f:
            line = f.read().strip()
            username, password = line.split(':', 1)
            return {username: password}
    except FileNotFoundError:
        print("WARNING: creds.txt not found. Creating default credentials.")
        with open('creds.txt', 'w') as f:
            f.write('admin:changeme123')
        return {'admin': 'changeme123'}


def login_required(f):
    """Decorator to require login for routes"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'logged_in' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function


def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def get_tv_list():
    """Get list of all TV directories"""
    tvs = []
    for item in os.listdir(SLIDES_DIR):
        item_path = os.path.join(SLIDES_DIR, item)
        if os.path.isdir(item_path) and item.startswith('TV_'):
            tvs.append(item)
    return sorted(tvs)


def get_tv_config(tv_id):
    """Load configuration for a specific TV"""
    config_path = os.path.join(SLIDES_DIR, tv_id, 'config.json')
    if os.path.exists(config_path):
        with open(config_path, 'r') as f:
            return json.load(f)
    return []


def save_tv_config(tv_id, config):
    """Save configuration for a specific TV"""
    tv_dir = os.path.join(SLIDES_DIR, tv_id)
    os.makedirs(tv_dir, exist_ok=True)
    config_path = os.path.join(tv_dir, 'config.json')
    with open(config_path, 'w') as f:
        json.dump(config, f, indent=2)


def convert_pptx_to_images(pptx_path, tv_id):
    """
    Convert PowerPoint presentation to JPG images
    Returns list of slide configs with filenames and durations
    """
    prs = Presentation(pptx_path)
    slides_config = []
    tv_dir = os.path.join(SLIDES_DIR, tv_id)

    # Get existing slides to determine next number
    existing_files = [f for f in os.listdir(tv_dir) if f.endswith('.jpg')]
    next_num = len(existing_files) + 1

    for idx, slide in enumerate(prs.slides):
        # Try to get slide duration from timing
        duration = DEFAULT_SLIDE_DURATION
        try:
            # Check if slide has timing set (in milliseconds)
            if hasattr(slide, 'timing') and slide.timing:
                duration = slide.timing / 1000.0
        except:
            pass

        # Generate filename
        filename = f"{next_num + idx}.jpg"
        filepath = os.path.join(tv_dir, filename)

        # Export slide as image
        # Note: python-pptx doesn't directly support rendering slides to images
        # We'll use a workaround by saving as image through Pillow
        # For production, consider using LibreOffice headless or similar

        # For now, we'll create a placeholder approach
        # In a real scenario, you'd use: unoconv, libreoffice --headless, or similar
        # This simplified version assumes slides are already images or we use a conversion service

        try:
            # Try to extract images from slide shapes
            image_found = False
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image_stream = io.BytesIO(shape.image.blob)
                    img = Image.open(image_stream)

                    # Resize to target resolution maintaining aspect ratio
                    img = resize_image(img, TARGET_WIDTH, TARGET_HEIGHT)
                    img = img.convert('RGB')
                    img.save(filepath, 'JPEG', quality=90)
                    image_found = True
                    break

            if not image_found:
                # Create a blank slide with text if no image found
                img = Image.new('RGB', (TARGET_WIDTH, TARGET_HEIGHT), color='white')
                img.save(filepath, 'JPEG', quality=90)
        except Exception as e:
            print(f"Error processing slide {idx}: {e}")
            # Create blank placeholder
            img = Image.new('RGB', (TARGET_WIDTH, TARGET_HEIGHT), color='lightgray')
            img.save(filepath, 'JPEG', quality=90)

        slides_config.append({
            'filename': filename,
            'duration_seconds': duration
        })

    return slides_config


def resize_image(img, target_width, target_height):
    """Resize image to target dimensions while maintaining aspect ratio"""
    # Calculate aspect ratios
    img_ratio = img.width / img.height
    target_ratio = target_width / target_height

    if img_ratio > target_ratio:
        # Image is wider - fit to width
        new_width = target_width
        new_height = int(target_width / img_ratio)
    else:
        # Image is taller - fit to height
        new_height = target_height
        new_width = int(target_height * img_ratio)

    # Resize image
    img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

    # Create new image with target dimensions and paste resized image centered
    new_img = Image.new('RGB', (target_width, target_height), color='black')
    paste_x = (target_width - new_width) // 2
    paste_y = (target_height - new_height) // 2
    new_img.paste(img, (paste_x, paste_y))

    return new_img


@app.route('/')
def index():
    """Redirect to dashboard if logged in, otherwise to login"""
    if 'logged_in' in session:
        return redirect(url_for('dashboard'))
    return redirect(url_for('login'))


@app.route('/login', methods=['GET', 'POST'])
def login():
    """Login page"""
    if request.method == 'POST':
        credentials = load_credentials()
        username = request.form.get('username')
        password = request.form.get('password')

        if username in credentials and credentials[username] == password:
            session['logged_in'] = True
            session['username'] = username
            return redirect(url_for('dashboard'))
        else:
            return render_template('login.html', error='Invalid credentials')

    return render_template('login.html')


@app.route('/logout')
def logout():
    """Logout user"""
    session.clear()
    return redirect(url_for('login'))


@app.route('/dashboard')
@login_required
def dashboard():
    """Main dashboard showing all TVs"""
    tvs = get_tv_list()
    tv_data = []

    for tv_id in tvs:
        config = get_tv_config(tv_id)
        tv_data.append({
            'id': tv_id,
            'slide_count': len(config)
        })

    return render_template('dashboard.html', tvs=tv_data, max_tvs=MAX_TVS)


@app.route('/tv/<tv_id>')
@login_required
def tv_management(tv_id):
    """Management page for specific TV"""
    if not tv_id.startswith('TV_'):
        return "Invalid TV ID", 400

    # Create TV directory if it doesn't exist
    tv_dir = os.path.join(SLIDES_DIR, tv_id)
    os.makedirs(tv_dir, exist_ok=True)

    config = get_tv_config(tv_id)
    return render_template('tv_management.html', tv_id=tv_id, slides=config)


@app.route('/display/<tv_id>')
def display(tv_id):
    """Fullscreen slideshow display for TV (no auth required)"""
    if not tv_id.startswith('TV_'):
        return "Invalid TV ID", 400

    config = get_tv_config(tv_id)
    return render_template('display.html', tv_id=tv_id, slides=config)


@app.route('/api/check/<tv_id>')
def api_check(tv_id):
    """API endpoint to check for config updates"""
    if not tv_id.startswith('TV_'):
        return jsonify({'error': 'Invalid TV ID'}), 400

    config = get_tv_config(tv_id)
    return jsonify(config)


@app.route('/api/upload_pptx/<tv_id>', methods=['POST'])
@login_required
def api_upload_pptx(tv_id):
    """Upload and convert PowerPoint presentation"""
    if not tv_id.startswith('TV_'):
        return jsonify({'error': 'Invalid TV ID'}), 400

    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    if not file.filename.endswith('.pptx'):
        return jsonify({'error': 'Only .pptx files allowed'}), 400

    # Save uploaded file temporarily
    temp_path = os.path.join(SLIDES_DIR, tv_id, 'temp_upload.pptx')
    file.save(temp_path)

    try:
        # Convert PowerPoint to images
        new_slides = convert_pptx_to_images(temp_path, tv_id)

        # Add to existing config
        config = get_tv_config(tv_id)
        config.extend(new_slides)
        save_tv_config(tv_id, config)

        # Remove temp file
        os.remove(temp_path)

        return jsonify({'success': True, 'slides_added': len(new_slides)})
    except Exception as e:
        if os.path.exists(temp_path):
            os.remove(temp_path)
        return jsonify({'error': str(e)}), 500


@app.route('/api/upload_image/<tv_id>', methods=['POST'])
@login_required
def api_upload_image(tv_id):
    """Upload individual image slide"""
    if not tv_id.startswith('TV_'):
        return jsonify({'error': 'Invalid TV ID'}), 400

    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type'}), 400

    # Generate filename
    tv_dir = os.path.join(SLIDES_DIR, tv_id)
    os.makedirs(tv_dir, exist_ok=True)

    existing_files = [f for f in os.listdir(tv_dir) if f.endswith('.jpg')]
    next_num = len(existing_files) + 1
    filename = f"{next_num}.jpg"
    filepath = os.path.join(tv_dir, filename)

    # Process and save image
    try:
        img = Image.open(file.stream)
        img = resize_image(img, TARGET_WIDTH, TARGET_HEIGHT)
        img = img.convert('RGB')
        img.save(filepath, 'JPEG', quality=90)

        # Add to config
        config = get_tv_config(tv_id)
        duration = float(request.form.get('duration', DEFAULT_SLIDE_DURATION))
        config.append({
            'filename': filename,
            'duration_seconds': duration
        })
        save_tv_config(tv_id, config)

        return jsonify({'success': True, 'filename': filename})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/reorder/<tv_id>', methods=['POST'])
@login_required
def api_reorder(tv_id):
    """Reorder slides"""
    if not tv_id.startswith('TV_'):
        return jsonify({'error': 'Invalid TV ID'}), 400

    new_order = request.json.get('order', [])
    config = get_tv_config(tv_id)

    # Create new config based on order
    new_config = []
    for filename in new_order:
        # Find slide in current config
        for slide in config:
            if slide['filename'] == filename:
                new_config.append(slide)
                break

    save_tv_config(tv_id, new_config)
    return jsonify({'success': True})


@app.route('/api/update_duration/<tv_id>', methods=['POST'])
@login_required
def api_update_duration(tv_id):
    """Update slide duration"""
    if not tv_id.startswith('TV_'):
        return jsonify({'error': 'Invalid TV ID'}), 400

    filename = request.json.get('filename')
    duration = float(request.json.get('duration', DEFAULT_SLIDE_DURATION))

    config = get_tv_config(tv_id)
    for slide in config:
        if slide['filename'] == filename:
            slide['duration_seconds'] = duration
            break

    save_tv_config(tv_id, config)
    return jsonify({'success': True})


@app.route('/api/delete_slide/<tv_id>', methods=['POST'])
@login_required
def api_delete_slide(tv_id):
    """Delete a slide"""
    if not tv_id.startswith('TV_'):
        return jsonify({'error': 'Invalid TV ID'}), 400

    filename = request.json.get('filename')

    # Remove from config
    config = get_tv_config(tv_id)
    config = [s for s in config if s['filename'] != filename]
    save_tv_config(tv_id, config)

    # Delete file
    filepath = os.path.join(SLIDES_DIR, tv_id, filename)
    if os.path.exists(filepath):
        os.remove(filepath)

    return jsonify({'success': True})


@app.route('/api/create_tv', methods=['POST'])
@login_required
def api_create_tv():
    """Create a new TV"""
    tv_list = get_tv_list()
    if len(tv_list) >= MAX_TVS:
        return jsonify({'error': f'Maximum {MAX_TVS} TVs allowed'}), 400

    # Find next available TV number
    existing_nums = []
    for tv in tv_list:
        try:
            num = int(tv.split('_')[1])
            existing_nums.append(num)
        except:
            pass

    next_num = 1
    while next_num in existing_nums:
        next_num += 1

    new_tv_id = f"TV_{next_num:03d}"
    tv_dir = os.path.join(SLIDES_DIR, new_tv_id)
    os.makedirs(tv_dir, exist_ok=True)
    save_tv_config(new_tv_id, [])

    return jsonify({'success': True, 'tv_id': new_tv_id})


@app.route('/slides/<tv_id>/<filename>')
def serve_slide(tv_id, filename):
    """Serve slide images"""
    tv_dir = os.path.join(SLIDES_DIR, tv_id)
    return send_from_directory(tv_dir, filename)


if __name__ == '__main__':
    # Run Flask development server
    # For production, use a proper WSGI server
    app.run(host='0.0.0.0', port=5000, debug=True)
